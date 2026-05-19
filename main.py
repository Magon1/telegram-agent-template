import os
import re
import json
import io
import time
import asyncio
import base64
from datetime import datetime, timedelta, timezone
from collections import deque
from contextlib import asynccontextmanager
from fastapi import FastAPI, Request
from fastapi.responses import RedirectResponse, HTMLResponse
import httpx
from anthropic import AsyncAnthropic
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import Flow
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from google.auth.transport.requests import Request as GoogleRequest
from pinecone import Pinecone
import voyageai
from pypdf import PdfReader
from pptx import Presentation
from docx import Document
from telethon import TelegramClient
from telethon.sessions import StringSession
from apscheduler.schedulers.asyncio import AsyncIOScheduler
from apscheduler.triggers.cron import CronTrigger
import pytz


# ==================== Lifespan ====================

_scheduler = None

@asynccontextmanager
async def lifespan(app: FastAPI):
    global _scheduler
    _scheduler = AsyncIOScheduler()
    seoul_tz = pytz.timezone('Asia/Seoul')
    _scheduler.add_job(
        send_morning_briefing,
        CronTrigger(hour=7, minute=0, timezone=seoul_tz),
        id='morning_briefing',
        replace_existing=True
    )
    _scheduler.start()
    print("[scheduler] Morning briefing: 07:00 KST daily")
    yield
    if _scheduler:
        _scheduler.shutdown()


app = FastAPI(lifespan=lifespan)


# ==================== 환경변수 ====================

_clean = lambda s: re.sub(r'\s+', '', s or '')

# 필수: Telegram + Anthropic
TELEGRAM_BOT_TOKEN = _clean(os.getenv("TELEGRAM_BOT_TOKEN"))
AUTHORIZED_CHAT_ID = int(_clean(os.getenv("AUTHORIZED_CHAT_ID")) or "0")
ANTHROPIC_API_KEY = _clean(os.getenv("ANTHROPIC_API_KEY"))

# 사용자 개인화
USER_NAME = os.getenv("USER_NAME", "사용자")
USER_TITLE = os.getenv("USER_TITLE", "")
COMPANY_NAME = os.getenv("COMPANY_NAME", "회사")
COMPANY_DESCRIPTION = os.getenv("COMPANY_DESCRIPTION", "")
USER_INTERESTS = os.getenv("USER_INTERESTS", "")
KEY_PEOPLE = os.getenv("KEY_PEOPLE", "")

# Google (선택)
GOOGLE_CLIENT_ID = _clean(os.getenv("GOOGLE_CLIENT_ID"))
GOOGLE_CLIENT_SECRET = _clean(os.getenv("GOOGLE_CLIENT_SECRET"))
GOOGLE_REFRESH_TOKEN = _clean(os.getenv("GOOGLE_REFRESH_TOKEN"))
GOOGLE_DRIVE_KB_FOLDER_ID = _clean(os.getenv("GOOGLE_DRIVE_KB_FOLDER_ID"))

# Slack (선택)
SLACK_TOKEN = _clean(os.getenv("SLACK_BOT_TOKEN") or os.getenv("SLACK_USER_TOKEN"))

# Pinecone + Voyage (선택, RAG/Memory용)
PINECONE_API_KEY = _clean(os.getenv("PINECONE_API_KEY"))
PINECONE_INDEX_NAME = _clean(os.getenv("PINECONE_INDEX_NAME")) or "knowledge-base"
VOYAGE_API_KEY = _clean(os.getenv("VOYAGE_API_KEY"))

# Telethon (선택, DM 모니터링용)
TELEGRAM_API_ID = int(_clean(os.getenv("TELEGRAM_API_ID")) or "0")
TELEGRAM_API_HASH = _clean(os.getenv("TELEGRAM_API_HASH"))
TELEGRAM_PHONE = _clean(os.getenv("TELEGRAM_PHONE"))
TELETHON_SESSION = _clean(os.getenv("TELETHON_SESSION"))

# Railway URL (자동)
RAILWAY_URL = _clean(os.getenv("RAILWAY_PUBLIC_DOMAIN") or "")

REDIRECT_URI = f"https://{RAILWAY_URL}/auth/google/callback"
GOOGLE_SCOPES = [
    'https://www.googleapis.com/auth/calendar',
    'https://www.googleapis.com/auth/drive.readonly',
    'https://www.googleapis.com/auth/gmail.readonly',
]

claude = AsyncAnthropic(api_key=ANTHROPIC_API_KEY)
conversation_history = deque(maxlen=20)

pc = Pinecone(api_key=PINECONE_API_KEY) if PINECONE_API_KEY else None
pinecone_index = pc.Index(PINECONE_INDEX_NAME) if pc else None
voyage_client = voyageai.Client(api_key=VOYAGE_API_KEY) if VOYAGE_API_KEY else None

indexed_files_cache = set()
_dm_scan_running = False
_morning_briefing_running = False


# ==================== 팀 컨텍스트 (파일에서 로드) ====================

def load_team_context():
    try:
        with open('team_context.txt', 'r', encoding='utf-8') as f:
            return f.read().strip()
    except FileNotFoundError:
        return "[team_context.txt 파일을 만들어 팀/관계자 정보를 입력하세요]"

TEAM_CONTEXT = load_team_context()

PRIORITY_RULES = """
🔴 긴급 답장 필요 — 24시간 이상 답 없는 중요 메시지, 액션 필요
🟡 답장 대기 — 답변 필요한 중요정보, 팔로업
🟢 정보 — 알아만 두면 됨
⚪ 무시 가능 — 인사·잡담
"""


# ==================== 메모리 시스템 ====================

MEMORY_NAMESPACE = "memory"


def add_memory(text: str, source: str = "manual") -> dict:
    if not pinecone_index or not voyage_client:
        return {"error": "Pinecone/Voyage 미설정"}
    try:
        memory_id = f"mem_{int(time.time() * 1000)}"
        emb = voyage_client.embed(texts=[text], model="voyage-3", input_type="document").embeddings[0]
        pinecone_index.upsert(
            vectors=[{
                "id": memory_id,
                "values": emb,
                "metadata": {
                    "text": text[:2000],
                    "created_at": datetime.now().isoformat(),
                    "source": source
                }
            }],
            namespace=MEMORY_NAMESPACE
        )
        return {"success": True, "memory_id": memory_id, "text": text}
    except Exception as e:
        return {"error": str(e)}


def search_memory(query: str, top_k: int = 5, min_score: float = 0.5) -> dict:
    if not pinecone_index or not voyage_client:
        return {"error": "Pinecone/Voyage 미설정"}
    try:
        query_emb = voyage_client.embed(texts=[query], model="voyage-3", input_type="query").embeddings[0]
        results = pinecone_index.query(
            vector=query_emb,
            top_k=top_k,
            include_metadata=True,
            namespace=MEMORY_NAMESPACE
        )
        relevant = [m for m in results.matches if m.score >= min_score]
        return {
            "count": len(relevant),
            "memories": [{
                "id": m.id,
                "text": m.metadata.get("text", ""),
                "created_at": m.metadata.get("created_at", ""),
                "score": round(m.score, 3)
            } for m in relevant]
        }
    except Exception as e:
        return {"error": str(e)}


def list_memories(limit: int = 30) -> dict:
    if not pinecone_index:
        return {"error": "Pinecone 미설정"}
    try:
        zero_vec = [0.0] * 1024
        results = pinecone_index.query(
            vector=zero_vec,
            top_k=limit,
            include_metadata=True,
            namespace=MEMORY_NAMESPACE
        )
        memories = [{
            "id": m.id,
            "text": m.metadata.get("text", ""),
            "created_at": m.metadata.get("created_at", "")
        } for m in results.matches]
        memories.sort(key=lambda x: x.get("created_at", ""), reverse=True)
        return {"count": len(memories), "memories": memories}
    except Exception as e:
        return {"error": str(e)}


def delete_memories_by_keyword(keyword: str) -> dict:
    if not pinecone_index or not voyage_client:
        return {"error": "Pinecone/Voyage 미설정"}
    try:
        search_result = search_memory(keyword, top_k=20, min_score=0.0)
        if "error" in search_result:
            return search_result
        matches = [m for m in search_result['memories'] if keyword.lower() in m['text'].lower()]
        if not matches:
            return {"error": f"'{keyword}' 매칭 기억 없음"}
        ids = [m['id'] for m in matches]
        pinecone_index.delete(ids=ids, namespace=MEMORY_NAMESPACE)
        return {
            "success": True,
            "deleted_count": len(ids),
            "deleted_memories": [m['text'][:100] for m in matches]
        }
    except Exception as e:
        return {"error": str(e)}


def memory_stats() -> dict:
    if not pinecone_index:
        return {"count": 0}
    try:
        stats = pinecone_index.describe_index_stats()
        ns = stats.get('namespaces', {}).get(MEMORY_NAMESPACE, {})
        return {"count": ns.get('vector_count', 0)}
    except Exception:
        return {"count": 0}


# ==================== 모델 라우터 ====================

def select_model(user_message: str) -> str:
    msg = user_message.lower().strip()
    if any(msg.startswith(p) for p in ("/opus", "/think", "/deep", "/깊게")):
        return "claude-opus-4-7"
    if any(msg.startswith(p) for p in ("/sonnet", "/balanced")):
        return "claude-sonnet-4-6"
    if any(msg.startswith(p) for p in ("/haiku", "/quick", "/fast", "/빠르게")):
        return "claude-haiku-4-5-20251001"
    sonnet_keywords = [
        "전략", "기획", "방향성", "로드맵", "방안", "대안",
        "분석", "검토", "비교", "장단점", "리스크", "전망", "예측",
        "어떻게 생각", "조언", "추천", "솔직하게", "의견", "판단",
        "어떡하지", "어떻게 해야", "고민",
        "초안", "보고서", "제안서", "정리해",
        "bm", "비즈니스 모델", "ir", "투자", "vc", "딜", "deal",
        "valuation", "밸류에이션", "왜", "이유", "근거"
    ]
    for kw in sonnet_keywords:
        if kw in msg:
            return "claude-sonnet-4-6"
    if len(user_message) > 200:
        return "claude-sonnet-4-6"
    return "claude-haiku-4-5-20251001"


# ==================== 텍스트 추출 / RAG ====================

def extract_text_from_pdf(file_bytes: bytes) -> list:
    reader = PdfReader(io.BytesIO(file_bytes))
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            pages.append({"page": i + 1, "text": text})
    return pages


def extract_text_from_pptx(file_bytes: bytes) -> list:
    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        combined = "\n".join(texts).strip()
        if combined:
            slides.append({"page": i + 1, "text": combined})
    return slides


def extract_text_from_docx(file_bytes: bytes) -> list:
    doc = Document(io.BytesIO(file_bytes))
    paragraphs = []
    current_chunk = []
    chunk_idx = 1
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            current_chunk.append(text)
            if len(" ".join(current_chunk)) > 800:
                paragraphs.append({"page": chunk_idx, "text": "\n".join(current_chunk)})
                current_chunk = []
                chunk_idx += 1
    if current_chunk:
        paragraphs.append({"page": chunk_idx, "text": "\n".join(current_chunk)})
    return paragraphs


def extract_text(filename: str, file_bytes: bytes) -> list:
    name_lower = filename.lower()
    if name_lower.endswith(".pdf"):
        return extract_text_from_pdf(file_bytes)
    elif name_lower.endswith(".pptx"):
        return extract_text_from_pptx(file_bytes)
    elif name_lower.endswith(".docx"):
        return extract_text_from_docx(file_bytes)
    elif name_lower.endswith(".txt") or name_lower.endswith(".md"):
        text = file_bytes.decode('utf-8', errors='ignore')
        chunks = [text[i:i+1000] for i in range(0, len(text), 1000)]
        return [{"page": i+1, "text": c} for i, c in enumerate(chunks)]
    return []


def chunk_text(text: str, max_chars: int = 1500) -> list:
    if len(text) <= max_chars:
        return [text]
    chunks = []
    current = ""
    for sentence in re.split(r'(?<=[.!?。!?])\s+', text):
        if len(current) + len(sentence) > max_chars:
            if current:
                chunks.append(current.strip())
            current = sentence
        else:
            current += " " + sentence
    if current.strip():
        chunks.append(current.strip())
    return chunks


def embed_texts(texts: list) -> list:
    if not voyage_client:
        return []
    result = voyage_client.embed(texts=texts, model="voyage-3", input_type="document")
    return result.embeddings


async def index_document(filename: str, file_bytes: bytes, source: str = "telegram") -> dict:
    if not pinecone_index or not voyage_client:
        return {"error": "Pinecone 또는 Voyage 미설정"}
    pages = extract_text(filename, file_bytes)
    if not pages:
        return {"error": f"지원 안 하거나 빈 파일: {filename}"}
    vectors_to_upsert = []
    chunk_count = 0
    for page_info in pages:
        chunks = chunk_text(page_info["text"], max_chars=1500)
        if not chunks:
            continue
        embeddings = embed_texts(chunks)
        for i, (chunk, emb) in enumerate(zip(chunks, embeddings)):
            vectors_to_upsert.append({
                "id": f"{filename}::page{page_info['page']}::chunk{i}",
                "values": emb,
                "metadata": {
                    "filename": filename, "page": page_info['page'], "text": chunk[:2000],
                    "source": source, "indexed_at": datetime.now().isoformat()
                }
            })
            chunk_count += 1
    for i in range(0, len(vectors_to_upsert), 100):
        pinecone_index.upsert(vectors=vectors_to_upsert[i:i+100])
    return {"success": True, "filename": filename, "pages": len(pages), "chunks": chunk_count}


def search_knowledge_base(query: str, top_k: int = 5) -> dict:
    if not pinecone_index or not voyage_client:
        return {"error": "Pinecone 또는 Voyage 미설정"}
    query_emb = voyage_client.embed(texts=[query], model="voyage-3", input_type="query").embeddings[0]
    results = pinecone_index.query(vector=query_emb, top_k=top_k, include_metadata=True)
    return {"count": len(results.matches), "matches": [{
        "filename": m.metadata.get("filename", ""), "page": m.metadata.get("page", 0),
        "score": round(m.score, 3), "text": m.metadata.get("text", ""),
        "source": m.metadata.get("source", "")
    } for m in results.matches]}


def list_indexed_documents() -> dict:
    if not pinecone_index:
        return {"error": "Pinecone 미설정"}
    stats = pinecone_index.describe_index_stats()
    return {"total_chunks": stats.get("total_vector_count", 0)}


# ==================== Telethon ====================

async def scan_telegram_messages(hours_back: int = 24) -> dict:
    if not all([TELEGRAM_API_ID, TELEGRAM_API_HASH, TELETHON_SESSION]):
        return {"error": "Telethon 환경변수 미설정"}
    client = TelegramClient(StringSession(TELETHON_SESSION), TELEGRAM_API_ID, TELEGRAM_API_HASH)
    try:
        await client.connect()
        if not await client.is_user_authorized():
            return {"error": "Telethon 세션 만료. 로컬 재인증 필요"}
        cutoff = datetime.now(timezone.utc) - timedelta(hours=hours_back)
        dialogs_data = []
        async for dialog in client.iter_dialogs(limit=80):
            if dialog.date and dialog.date < cutoff:
                continue
            entity = dialog.entity
            username = getattr(entity, 'username', None)
            entity_id = getattr(entity, 'id', None)
            is_linkable = dialog.is_channel
            messages = []
            try:
                async for msg in client.iter_messages(dialog.entity, limit=30):
                    if not msg.date or msg.date < cutoff:
                        break
                    if not msg.text:
                        continue
                    sender_name = "?"
                    if msg.sender:
                        sender_name = (
                            getattr(msg.sender, 'first_name', None) or
                            getattr(msg.sender, 'title', None) or
                            getattr(msg.sender, 'username', None) or
                            "Unknown"
                        )
                    msg_link = None
                    if is_linkable and msg.id:
                        if username:
                            msg_link = f"https://t.me/{username}/{msg.id}"
                        elif entity_id:
                            clean_id = abs(entity_id)
                            if str(clean_id).startswith('100'):
                                clean_id = int(str(clean_id)[3:])
                            msg_link = f"https://t.me/c/{clean_id}/{msg.id}"
                    messages.append({
                        'sender': f"{USER_NAME}(본인)" if msg.out else sender_name,
                        'text': (msg.text or '')[:400],
                        'date': msg.date.strftime('%m-%d %H:%M'),
                        'is_me': bool(msg.out),
                        'link': msg_link
                    })
            except Exception as e:
                print(f"[telethon msg error] {dialog.name}: {e}")
                continue
            if messages:
                if dialog.is_channel and not dialog.is_group:
                    dialog_type = "channel"
                elif dialog.is_group:
                    dialog_type = "group"
                else:
                    dialog_type = "dm"
                channel_link = f"https://t.me/{username}" if username else None
                dialogs_data.append({
                    'name': dialog.name or "Unknown",
                    'type': dialog_type,
                    'unread_count': dialog.unread_count,
                    'message_count': len(messages),
                    'messages': messages,
                    'channel_link': channel_link
                })
        return {
            'scanned_at': datetime.now().isoformat(),
            'hours_back': hours_back,
            'dialog_count': len(dialogs_data),
            'dialogs': dialogs_data
        }
    finally:
        await client.disconnect()


async def summarize_telegram_activity(hours_back: int = 24) -> dict:
    scan_result = await scan_telegram_messages(hours_back)
    if "error" in scan_result:
        return scan_result
    if scan_result['dialog_count'] == 0:
        return {"summary": f"최근 {hours_back}시간 텔레그램 활동 없음"}
    
    dialogs_summary = ""
    for d in scan_result['dialogs'][:40]:
        ch_link = f" [채널링크: {d['channel_link']}]" if d.get('channel_link') else ""
        last_status = ""
        if d['messages']:
            latest = d['messages'][0]
            who = f"✅{USER_NAME}(본인)" if latest['is_me'] else f"❗{latest['sender']}"
            last_status = f"\n[마지막 메시지 발신: {who} / 시각: {latest['date']}]"
        dialogs_summary += f"\n\n=== {d['name']} ({d['type']}, 메시지 {d['message_count']}개, 안 읽음 {d['unread_count']}){ch_link} ==={last_status}\n"
        for m in d['messages'][:20]:
            prefix = "[나]" if m['is_me'] else f"[{m['sender']}]"
            link_part = f" [msg_link: {m['link']}]" if m.get('link') else ""
            dialogs_summary += f"{m['date']} {prefix}: {m['text']}{link_part}\n"
    
    scan_time = datetime.now().strftime('%Y.%m.%d %H:%M')
    
    prompt = f"""다음은 {USER_NAME}님의 텔레그램 최근 {hours_back}시간 활동입니다.
스캔 시각: {scan_time} KST

{TEAM_CONTEXT}

{PRIORITY_RULES}

[메시지 데이터]
{dialogs_summary}

🚨🚨🚨 *답장 여부 판단 - 가장 중요한 규칙*

각 대화의 첫 줄 [마지막 메시지 발신: ~ / 시각: ~] 표기를 반드시 확인.

판단 절차:
1. 마지막 발신이 ✅{USER_NAME}(본인) → 이미 답장 완료
   - 🔴/🟡 절대 X. 정보 가치 있으면 🟢, 잡담이면 ⚪
2. 마지막 발신이 ❗상대방 → {USER_NAME}님 답장 아직 없음
   - 응답 요구 명확 + 24시간 경과 → 🔴
   - 응답 요구 있지만 시간 여유 → 🟡
   - 단순 정보/통보 → 🟢
   - 인사/잡담 → ⚪

분류 원칙:
1. "개인 DM"+"팀 그룹 대화"만 우선순위 분류
2. "정보 채널"(다수 모인 정보/잡담성)은 우선순위 X → '오늘의 인사이트'에만
3. ⚪ 무시 가능은 개인/팀 대화 중 잡담만 (정보 채널 절대 X)
4. 보안 위협 🔴 최상단

출력 포맷 (텔레그램 Markdown, 엄격):
- '---' '#' '##' '###' 헤더 X
- 섹션 제목은 *별표 양쪽* 볼드만
- 인사이트 근거 채널은 [채널명](msg_link) 마크다운 링크

출력 형식:

📋 *텔레그램 동향* — {scan_time}


*🔴 긴급 답장 필요*

- [상대/그룹명] 핵심 1줄
   → 명령조 액션


*🟡 답장 대기*

- [상대/그룹명] 핵심
   → 명령조 액션


*🟢 정보*

- [상대/그룹명] 핵심


*⚪ 무시 가능*: N개


*🔥 오늘의 인사이트*

① *제목*
   > 근거: [채널1](msg_link1), [채널2](msg_link2)
   > 1-2줄 요약 + 비즈니스 연결점

인사이트 최소 2-3개.

🚨 최종 검토: 🔴/🟡에 넣기 전 [마지막 메시지 발신: ✅{USER_NAME}] 이면 빼야 함"""
    
    try:
        response = await claude.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=4096,
            messages=[{"role": "user", "content": prompt}]
        )
        return {'summary': response.content[0].text, 'dialog_count': scan_result['dialog_count']}
    except Exception as e:
        return {"error": f"Claude 분류 에러: {str(e)[:200]}"}


# ==================== Google Calendar / Drive / Gmail ====================

def get_google_credentials():
    if not GOOGLE_REFRESH_TOKEN:
        return None
    creds = Credentials(
        token=None, refresh_token=GOOGLE_REFRESH_TOKEN,
        token_uri="https://oauth2.googleapis.com/token",
        client_id=GOOGLE_CLIENT_ID, client_secret=GOOGLE_CLIENT_SECRET,
        scopes=GOOGLE_SCOPES,
    )
    creds.refresh(GoogleRequest())
    return creds


def get_drive_service():
    creds = get_google_credentials()
    return build('drive', 'v3', credentials=creds) if creds else None


async def sync_drive_folder():
    if not GOOGLE_DRIVE_KB_FOLDER_ID:
        return {"error": "GOOGLE_DRIVE_KB_FOLDER_ID 미설정"}
    service = get_drive_service()
    if not service:
        return {"error": "Drive 인증 안 됨"}
    query = f"'{GOOGLE_DRIVE_KB_FOLDER_ID}' in parents and trashed=false"
    results = service.files().list(q=query, fields="files(id, name, mimeType, size, modifiedTime)", pageSize=100).execute()
    files = results.get('files', [])
    indexed, skipped, failed = [], [], []
    for f in files:
        file_id, filename = f['id'], f['name']
        cache_key = f"drive::{file_id}::{f.get('modifiedTime', '')}"
        if cache_key in indexed_files_cache:
            skipped.append(filename)
            continue
        if not filename.lower().endswith(('.pdf', '.pptx', '.docx', '.txt', '.md')):
            skipped.append(f"{filename} (지원 안 함)")
            continue
        try:
            request = service.files().get_media(fileId=file_id)
            file_bytes = io.BytesIO()
            downloader = MediaIoBaseDownload(file_bytes, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            result = await index_document(filename, file_bytes.getvalue(), source="drive")
            if result.get("success"):
                indexed.append(f"{filename} ({result['pages']}p, {result['chunks']}c)")
                indexed_files_cache.add(cache_key)
            else:
                failed.append(f"{filename}: {result.get('error')}")
        except Exception as e:
            failed.append(f"{filename}: {str(e)[:100]}")
    return {
        "indexed_count": len(indexed), "skipped_count": len(skipped), "failed_count": len(failed),
        "indexed": indexed, "failed": failed[:5]
    }


def list_calendars():
    creds = get_google_credentials()
    if not creds:
        return {"error": "Google 미인증"}
    service = build('calendar', 'v3', credentials=creds)
    cal_list = service.calendarList().list().execute()
    return {"count": len(cal_list.get('items', [])),
            "calendars": [{"id": c['id'], "name": c.get('summary'), "primary": c.get('primary', False)}
                          for c in cal_list.get('items', [])]}


def get_calendar_events(start_date: str, end_date: str, calendar_id: str = "primary"):
    creds = get_google_credentials()
    if not creds:
        return {"error": "Google 미인증"}
    service = build('calendar', 'v3', credentials=creds)
    try:
        events_result = service.events().list(
            calendarId=calendar_id,
            timeMin=f"{start_date}T00:00:00+09:00",
            timeMax=f"{end_date}T23:59:59+09:00",
            singleEvents=True, orderBy='startTime', maxResults=50,
        ).execute()
    except Exception as e:
        return {"error": f"캘린더 조회 실패: {str(e)[:200]}"}
    events = events_result.get('items', [])
    return {"calendar_id": calendar_id, "count": len(events), "events": [{
        "title": e.get('summary', '제목 없음'),
        "start": e['start'].get('dateTime', e['start'].get('date')),
        "end": e['end'].get('dateTime', e['end'].get('date')),
        "description": (e.get('description') or '')[:200],
        "location": e.get('location', ''),
    } for e in events]}


def create_calendar_event(title: str, start_datetime: str, end_datetime: str, description: str = ""):
    creds = get_google_credentials()
    if not creds:
        return {"error": "Google 미인증"}
    service = build('calendar', 'v3', credentials=creds)
    created = service.events().insert(calendarId='primary', body={
        'summary': title, 'description': description,
        'start': {'dateTime': start_datetime, 'timeZone': 'Asia/Seoul'},
        'end': {'dateTime': end_datetime, 'timeZone': 'Asia/Seoul'},
    }).execute()
    return {"success": True, "event_link": created.get('htmlLink'), "title": title}


def get_research_emails(hours_back: int = 24, max_results: int = 20) -> dict:
    creds = get_google_credentials()
    if not creds:
        return {"error": "Google 미인증"}
    service = build('gmail', 'v1', credentials=creds)
    cutoff_date = (datetime.now() - timedelta(hours=hours_back)).strftime('%Y/%m/%d')
    query = f"label:Research after:{cutoff_date}"
    try:
        results = service.users().messages().list(userId='me', q=query, maxResults=max_results).execute()
    except Exception as e:
        return {"error": f"Gmail 조회 실패: {str(e)[:200]}"}
    messages = results.get('messages', [])
    emails = []
    for msg in messages:
        try:
            msg_data = service.users().messages().get(
                userId='me', id=msg['id'], format='metadata',
                metadataHeaders=['Subject', 'From', 'Date']
            ).execute()
            headers = {h['name']: h['value'] for h in msg_data.get('payload', {}).get('headers', [])}
            emails.append({
                'subject': headers.get('Subject', '(제목 없음)'),
                'from': headers.get('From', '')[:80],
                'date': headers.get('Date', ''),
                'snippet': msg_data.get('snippet', '')[:500],
                'link': f"https://mail.google.com/mail/u/0/#inbox/{msg['id']}"
            })
        except Exception as e:
            print(f"[gmail msg error] {e}")
            continue
    return {'count': len(emails), 'emails': emails}


async def summarize_research_emails(hours_back: int = 24) -> str:
    result = get_research_emails(hours_back)
    if "error" in result:
        return f"❌ {result['error']}"
    if result['count'] == 0:
        return "_새 Research 메일 없음_"
    emails_text = ""
    for e in result['emails']:
        emails_text += f"\n[{e['from']}] {e['subject']}\n링크: {e['link']}\n{e['snippet']}\n"
    prompt = f"""다음은 {USER_NAME}님 Gmail의 Research 라벨 메일 {result['count']}건입니다.

{USER_NAME}님 관심사: {USER_INTERESTS}

[메일 데이터]
{emails_text}

출력 포맷:
- '#' '##' 헤더 X
- *별표 양쪽* 볼드만
- 중요도순 3-5건
- 형식: • *[발신자]* [제목](링크)
        > 핵심 1-2줄

핵심: 관심사 직결 우선, 숫자/날짜 보존, 제목 마크다운 링크"""
    try:
        response = await claude.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=2048,
            messages=[{"role": "user", "content": prompt}]
        )
        return response.content[0].text
    except Exception as e:
        return f"❌ Claude 에러: {str(e)[:200]}"


def summarize_today_calendar() -> str:
    today = datetime.now().strftime('%Y-%m-%d')
    result = get_calendar_events(today, today)
    if "error" in result:
        return f"❌ {result['error']}"
    if result['count'] == 0:
        return "_오늘 일정 없음_ 🎉"
    lines = []
    for e in result['events'][:10]:
        start = e.get('start', '')
        time_str = start.split('T')[1][:5] if 'T' in start else "하루종일"
        loc = f" @ {e['location']}" if e.get('location') else ""
        lines.append(f"• *{time_str}* — {e['title']}{loc}")
    return "\n".join(lines)


# ==================== 통합 아침 브리핑 ====================

async def send_morning_briefing():
    global _morning_briefing_running
    if _morning_briefing_running:
        await send_telegram_message(AUTHORIZED_CHAT_ID, "⚠️ 이미 브리핑 처리 중")
        return
    _morning_briefing_running = True
    try:
        await send_telegram_message(AUTHORIZED_CHAT_ID, "☀️ *아침 브리핑 생성 중...* (3-5분)")
        cal_summary = summarize_today_calendar()
        gmail_summary = await summarize_research_emails(hours_back=24)
        today = datetime.now().strftime('%Y.%m.%d (%A)')
        intro = (
            f"☀️ *Good morning, {USER_NAME}님*\n"
            f"_{today}_\n\n"
            f"📅 *오늘 일정*\n\n{cal_summary}\n\n"
            f"📧 *Research 메일*\n\n{gmail_summary}"
        )
        await send_telegram_message(AUTHORIZED_CHAT_ID, intro)
        tg_result = await summarize_telegram_activity(hours_back=24)
        if "error" in tg_result:
            await send_telegram_message(AUTHORIZED_CHAT_ID, f"💬 *텔레그램 동향*\n\n❌ {tg_result['error']}")
        else:
            await send_telegram_message(AUTHORIZED_CHAT_ID, tg_result.get('summary', '결과 없음'))
    except Exception as e:
        print(f"[morning briefing error] {e}")
        await send_telegram_message(AUTHORIZED_CHAT_ID, f"⚠️ 아침 브리핑 에러: {str(e)[:300]}")
    finally:
        _morning_briefing_running = False


# ==================== Slack ====================

async def slack_api_call(method: str, payload: dict = None):
    url = f"https://slack.com/api/{method}"
    headers = {"Authorization": f"Bearer {SLACK_TOKEN}", "Content-Type": "application/json; charset=utf-8"}
    async with httpx.AsyncClient(timeout=30) as client:
        if payload is None:
            response = await client.get(url, headers=headers)
        else:
            response = await client.post(url, headers=headers, json=payload)
        return response.json()


async def search_slack_channel(channel_name_or_keyword: str, limit: int = 30):
    if not SLACK_TOKEN:
        return {"error": "Slack 미설정"}
    list_result = await slack_api_call("conversations.list", None)
    channels = list_result.get("channels", [])
    matching = [c for c in channels if channel_name_or_keyword.lower() in c.get("name", "").lower()]
    if not matching:
        return {"error": "매칭 채널 없음", "available_channels": [c["name"] for c in channels[:20]]}
    results = {}
    for ch in matching[:3]:
        history = await slack_api_call("conversations.history", {"channel": ch["id"], "limit": limit})
        results[ch["name"]] = [{"user": m.get("user"), "text": m.get("text", "")[:500], "ts": m.get("ts")}
                                for m in history.get("messages", [])]
    return {"channels_found": list(results.keys()), "messages": results}


async def get_slack_channel_messages(channel_id: str, hours_ago: int = 24):
    if not SLACK_TOKEN:
        return {"error": "Slack 미설정"}
    result = await slack_api_call("conversations.history", {
        "channel": channel_id, "oldest": str(time.time() - hours_ago * 3600), "limit": 100
    })
    return {"count": len(result.get("messages", [])), "messages": [
        {"user": m.get("user"), "text": m.get("text", "")[:500], "ts": m.get("ts")}
        for m in result.get("messages", [])
    ]}


async def send_slack_message(channel_id: str, text: str):
    if not SLACK_TOKEN:
        return {"error": "Slack 미설정"}
    result = await slack_api_call("chat.postMessage", {"channel": channel_id, "text": text})
    return {"success": result.get("ok"), "error": result.get("error")}


async def list_slack_channels():
    if not SLACK_TOKEN:
        return {"error": "Slack 미설정"}
    result = await slack_api_call("conversations.list", None)
    return {"count": len(result.get("channels", [])), "channels": [
        {"id": c["id"], "name": c["name"], "is_member": c.get("is_member", False)}
        for c in result.get("channels", [])[:50]
    ]}


# ==================== Claude Tools ====================

CLAUDE_TOOLS = [
    {
        "name": "remember_fact",
        "description": "중요한 사실을 영구 기억에 저장. 사용자가 '기억해줘'/'저장해줘' 요청 시, 또는 대화 중 미래에도 필요할 정보 발견 시 자동 호출 (인물 정보, 일정, 결정 사항, 선호도, 비즈니스 컨텍스트 등). 한 문장으로 명확하게.",
        "input_schema": {"type": "object", "properties": {
            "fact": {"type": "string", "description": "기억할 사실 한 문장"}
        }, "required": ["fact"]}
    },
    {
        "name": "search_my_memory",
        "description": "사용자의 영구 기억에서 검색. '뭐였더라', '~ 정보 찾아줘', '예전에 알려준 ~' 같은 질문에 사용.",
        "input_schema": {"type": "object", "properties": {
            "query": {"type": "string"},
            "top_k": {"type": "integer", "default": 5}
        }, "required": ["query"]}
    },
    {
        "name": "scan_telegram_dms",
        "description": "텔레그램 DM/그룹/채널 최근 메시지 우선순위별 정리. '오늘 답장할 거?', '특정 팀 최근?' 등. 1-3분.",
        "input_schema": {"type": "object", "properties": {
            "hours_back": {"type": "integer", "default": 24}
        }}
    },
    {
        "name": "get_research_emails",
        "description": "Gmail Research 라벨 메일 조회. '리서치 메일', '뉴스레터'에 사용.",
        "input_schema": {"type": "object", "properties": {
            "hours_back": {"type": "integer", "default": 24}
        }}
    },
    {
        "name": "search_knowledge_base",
        "description": "업로드한 PDF/덱/딜자료 검색",
        "input_schema": {"type": "object", "properties": {
            "query": {"type": "string"},
            "top_k": {"type": "integer", "default": 5}
        }, "required": ["query"]}
    },
    {"name": "list_indexed_documents", "description": "지식 베이스 통계", "input_schema": {"type": "object", "properties": {}}},
    {"name": "sync_drive_folder", "description": "Drive KB 폴더 자동 인덱싱", "input_schema": {"type": "object", "properties": {}}},
    {"name": "list_calendars", "description": "캘린더 목록", "input_schema": {"type": "object", "properties": {}}},
    {
        "name": "get_calendar_events",
        "description": "캘린더 일정 조회",
        "input_schema": {"type": "object", "properties": {
            "start_date": {"type": "string"},
            "end_date": {"type": "string"},
            "calendar_id": {"type": "string", "default": "primary"}
        }, "required": ["start_date", "end_date"]}
    },
    {
        "name": "create_calendar_event",
        "description": "일정 추가",
        "input_schema": {"type": "object", "properties": {
            "title": {"type": "string"},
            "start_datetime": {"type": "string"},
            "end_datetime": {"type": "string"},
            "description": {"type": "string"}
        }, "required": ["title", "start_datetime", "end_datetime"]}
    },
    {"name": "list_slack_channels", "description": "슬랙 채널 목록", "input_schema": {"type": "object", "properties": {}}},
    {
        "name": "search_slack_channel",
        "description": "키워드로 슬랙 검색",
        "input_schema": {"type": "object", "properties": {
            "channel_name_or_keyword": {"type": "string"},
            "limit": {"type": "integer", "default": 30}
        }, "required": ["channel_name_or_keyword"]}
    },
    {
        "name": "get_slack_channel_messages",
        "description": "슬랙 채널 N시간 메시지",
        "input_schema": {"type": "object", "properties": {
            "channel_id": {"type": "string"},
            "hours_ago": {"type": "integer", "default": 24}
        }, "required": ["channel_id"]}
    },
    {
        "name": "send_slack_message",
        "description": "슬랙 발송 (명시적 요청 시만)",
        "input_schema": {"type": "object", "properties": {
            "channel_id": {"type": "string"},
            "text": {"type": "string"}
        }, "required": ["channel_id", "text"]}
    }
]


async def execute_tool(tool_name: str, tool_input: dict):
    try:
        if tool_name == "remember_fact":
            return add_memory(tool_input.get("fact", ""), source="auto")
        elif tool_name == "search_my_memory":
            return search_memory(**tool_input)
        elif tool_name == "scan_telegram_dms":
            return await summarize_telegram_activity(**tool_input)
        elif tool_name == "get_research_emails":
            return get_research_emails(**tool_input)
        elif tool_name == "search_knowledge_base":
            return search_knowledge_base(**tool_input)
        elif tool_name == "list_indexed_documents":
            return list_indexed_documents()
        elif tool_name == "sync_drive_folder":
            return await sync_drive_folder()
        elif tool_name == "list_calendars":
            return list_calendars()
        elif tool_name == "get_calendar_events":
            return get_calendar_events(**tool_input)
        elif tool_name == "create_calendar_event":
            return create_calendar_event(**tool_input)
        elif tool_name == "list_slack_channels":
            return await list_slack_channels()
        elif tool_name == "search_slack_channel":
            return await search_slack_channel(**tool_input)
        elif tool_name == "get_slack_channel_messages":
            return await get_slack_channel_messages(**tool_input)
        elif tool_name == "send_slack_message":
            return await send_slack_message(**tool_input)
        else:
            return {"error": f"알 수 없는 도구: {tool_name}"}
    except Exception as e:
        return {"error": str(e)}


SYSTEM_PROMPT = f"""당신은 {USER_TITLE} {USER_NAME}님의 개인 비서 AI입니다.

오늘 날짜: {datetime.now().strftime('%Y-%m-%d (%A)')}
타임존: Asia/Seoul

회사: {COMPANY_NAME}
{COMPANY_DESCRIPTION}

주요 인물: {KEY_PEOPLE}
관심사: {USER_INTERESTS}

{TEAM_CONTEXT}

도구:
[기억] remember_fact (영구 저장), search_my_memory (기억 검색)
[텔레그램 DM] scan_telegram_dms
[Gmail] get_research_emails
[지식 베이스] search_knowledge_base, list_indexed_documents, sync_drive_folder
[캘린더] list_calendars, get_calendar_events, create_calendar_event
[슬랙] list_slack_channels, search_slack_channel, get_slack_channel_messages, send_slack_message

원칙:
- 대화 중 미래에도 필요한 사실 (인물·역할, 일정, 결정, 선호도, 비즈니스 컨텍스트) 발견 시 → remember_fact 자동 호출
- '뭐였더라'/'예전에 알려준 ~'/'그거 기억나?' → search_my_memory 먼저
- '텔레그램 정리'/'답장 대기' → scan_telegram_dms
- '리서치 메일'/'뉴스레터' → get_research_emails
- 회사 내부 문서 → search_knowledge_base 먼저
- 슬랙 발송은 명시적 요청 시만
- 간단한 인사·잡담은 도구 호출 없이 즉답

스타일: 한국어, 존댓말, 간결, 정확. 모르면 모른다고.
"""


async def get_claude_response(user_message: str) -> str:
    # 관련 기억 자동 검색
    relevant_memories_text = ""
    try:
        mem_result = search_memory(user_message, top_k=5, min_score=0.5)
        if mem_result.get('count', 0) > 0:
            relevant_memories_text = "\n\n[자동 회상된 관련 기억]\n"
            for m in mem_result['memories']:
                relevant_memories_text += f"- {m['text']}\n"
    except Exception as e:
        print(f"[memory search error] {e}")
    
    system_with_memories = SYSTEM_PROMPT + relevant_memories_text
    
    conversation_history.append({"role": "user", "content": user_message})
    model = select_model(user_message)
    print(f"[model] {model} for: {user_message[:50]}")
    
    for iteration in range(10):
        try:
            response = await asyncio.wait_for(
                claude.messages.create(
                    model=model, max_tokens=4096,
                    system=system_with_memories,
                    tools=CLAUDE_TOOLS,
                    messages=list(conversation_history),
                ),
                timeout=60.0
            )
        except asyncio.TimeoutError:
            print(f"[claude timeout] iteration {iteration}")
            if conversation_history and conversation_history[-1].get("role") == "user":
                conversation_history.pop()
            return "⚠️ Claude 60초 초과. /reset 후 재시도해주세요."
        except Exception as e:
            print(f"[claude error] {e}")
            if conversation_history and conversation_history[-1].get("role") == "user":
                conversation_history.pop()
            return f"⚠️ Claude 에러: {str(e)[:200]}\n\n/reset 후 재시도."
        
        conversation_history.append({
            "role": "assistant",
            "content": [block.model_dump() for block in response.content]
        })
        
        if response.stop_reason == "tool_use":
            tool_results = []
            for block in response.content:
                if block.type == "tool_use":
                    print(f"[tool] {block.name}")
                    try:
                        tool_timeout = 300.0 if block.name in ("scan_telegram_dms", "sync_drive_folder") else 60.0
                        result = await asyncio.wait_for(
                            execute_tool(block.name, dict(block.input)),
                            timeout=tool_timeout
                        )
                    except asyncio.TimeoutError:
                        result = {"error": f"{block.name} 타임아웃 ({int(tool_timeout)}초)"}
                    except Exception as e:
                        result = {"error": str(e)}
                    tool_results.append({
                        "type": "tool_result",
                        "tool_use_id": block.id,
                        "content": json.dumps(result, ensure_ascii=False)[:12000]
                    })
            conversation_history.append({"role": "user", "content": tool_results})
        else:
            text_blocks = [b.text for b in response.content if b.type == "text"]
            return "\n".join(text_blocks) if text_blocks else "응답 비어있음 🤔"
    
    return "⚠️ 도구 호출 너무 많음. /reset 후 재시도."


# ==================== Telegram ====================

async def send_telegram_message(chat_id: int, text: str):
    url = f"https://api.telegram.org/bot{TELEGRAM_BOT_TOKEN}/sendMessage"
    if len(text) > 4000:
        chunks = [text[i:i+4000] for i in range(0, len(text), 4000)]
        async with httpx.AsyncClient(timeout=60) as client:
            for chunk in chunks:
                try:
                    await client.post(url, json={"chat_id": chat_id, "text": chunk, "parse_mode": "Markdown"})
                except Exception:
                    await client.post(url, json={"chat_id": chat_id, "text": chunk})
        return
    async with httpx.AsyncClient(timeout=60) as client:
        try:
            await client.post(url, json={"chat_id": chat_id, "text": text, "parse_mode": "Markdown"})
        except Exception:
            await client.post(url, json={"chat_id": chat_id, "text": text})


async def download_telegram_file(file_id: str) -> tuple:
    async with httpx.AsyncClient(timeout=120) as client:
        info = await client.get(
            f"https://api.telegram.org/bot{TELEGRAM_BOT_TOKEN}/getFile",
            params={"file_id": file_id}
        )
        info_json = info.json()
        if not info_json.get("ok"):
            raise Exception(f"파일 정보 조회 실패: {info_json}")
        file_path = info_json["result"]["file_path"]
        resp = await client.get(f"https://api.telegram.org/file/bot{TELEGRAM_BOT_TOKEN}/{file_path}")
        resp.raise_for_status()
        return file_path.split("/")[-1], resp.content


async def _run_dm_scan_task(chat_id: int, hours: int):
    global _dm_scan_running
    _dm_scan_running = True
    try:
        result = await summarize_telegram_activity(hours_back=hours)
        if "error" in result:
            await send_telegram_message(chat_id, f"❌ {result['error']}")
        else:
            await send_telegram_message(chat_id, result.get('summary', '결과 없음'))
    except Exception as e:
        await send_telegram_message(chat_id, f"❌ 처리 에러: {str(e)[:200]}")
    finally:
        _dm_scan_running = False


@app.post("/webhook")
async def telegram_webhook(request: Request):
    data = await request.json()
    if "message" not in data:
        return {"ok": True}
    
    message = data["message"]
    chat_id = message["chat"]["id"]
    if chat_id != AUTHORIZED_CHAT_ID:
        return {"ok": True}
    
    # 파일 업로드
    if "document" in message:
        doc = message["document"]
        file_id = doc["file_id"]
        filename = doc.get("file_name", "unknown")
        file_size_mb = doc.get("file_size", 0) / 1024 / 1024
        if file_size_mb > 20:
            await send_telegram_message(chat_id,
                f"⚠️ *{filename}* ({file_size_mb:.1f}MB) 20MB 초과\n💡 Drive 폴더 → /sync")
            return {"ok": True}
        await send_telegram_message(chat_id, f"📥 *{filename}* 받았어요. 처리 중...")
        try:
            _, file_bytes = await download_telegram_file(file_id)
            result = await index_document(filename, file_bytes, source="telegram")
            if result.get("success"):
                await send_telegram_message(chat_id, f"✅ *{filename}*\n📄 {result['pages']}p → {result['chunks']}c")
            else:
                await send_telegram_message(chat_id, f"❌ 실패: {result.get('error')}")
        except Exception as e:
            await send_telegram_message(chat_id, f"❌ 처리 에러: {str(e)[:200]}")
        return {"ok": True}
    
    text = message.get("text", "")
    
    # /remember
    if text.strip().startswith("/remember "):
        fact = text.strip()[len("/remember "):].strip()
        if not fact:
            await send_telegram_message(chat_id, "❌ 기억할 내용 필요\n예: `/remember Backpack 매니저는 Trish`")
            return {"ok": True}
        result = add_memory(fact, source="manual")
        if result.get("success"):
            await send_telegram_message(chat_id, f"✅ 기억했어요\n_{fact}_")
        else:
            await send_telegram_message(chat_id, f"❌ {result.get('error')}")
        return {"ok": True}
    
    # /memories
    if text.strip() == "/memories":
        result = list_memories(limit=30)
        if "error" in result:
            await send_telegram_message(chat_id, f"❌ {result['error']}")
        elif result['count'] == 0:
            await send_telegram_message(chat_id, "_저장된 기억 없음_\n\n`/remember <내용>` 으로 추가하세요")
        else:
            msg = f"🧠 *저장된 기억* ({result['count']}개)\n\n"
            for i, m in enumerate(result['memories'][:30], 1):
                msg += f"{i}. {m['text']}\n"
            await send_telegram_message(chat_id, msg)
        return {"ok": True}
    
    # /forget
    if text.strip().startswith("/forget "):
        pattern = text.strip()[len("/forget "):].strip()
        if not pattern:
            await send_telegram_message(chat_id, "❌ 키워드 필요\n예: `/forget Backpack`")
            return {"ok": True}
        result = delete_memories_by_keyword(pattern)
        if result.get("success"):
            msg = f"🗑 {result['deleted_count']}개 기억 삭제됨"
            for t in result.get('deleted_memories', []):
                msg += f"\n• {t}"
            await send_telegram_message(chat_id, msg)
        else:
            await send_telegram_message(chat_id, f"❌ {result.get('error')}")
        return {"ok": True}
    
    # /morning
    if text.strip() == "/morning":
        if _morning_briefing_running:
            await send_telegram_message(chat_id, "⚠️ 이미 브리핑 처리 중")
            return {"ok": True}
        asyncio.create_task(send_morning_briefing())
        return {"ok": True}
    
    # /dm
    if text.strip().startswith("/dm"):
        parts = text.strip().split()
        hours = 24
        if len(parts) > 1 and parts[1].isdigit():
            hours = int(parts[1])
        if _dm_scan_running:
            await send_telegram_message(chat_id, "⚠️ 이미 스캔 진행 중")
            return {"ok": True}
        await send_telegram_message(chat_id, f"📡 텔레그램 *{hours}시간* 스캔 시작...")
        asyncio.create_task(_run_dm_scan_task(chat_id, hours))
        return {"ok": True}
    
    if text.strip() == "/reset":
        conversation_history.clear()
        await send_telegram_message(chat_id, "🔄 대화 기억 초기화 (장기 기억은 그대로)")
        return {"ok": True}
    
    if text.strip() == "/sync":
        await send_telegram_message(chat_id, "🔄 Drive 폴더 스캔 중...")
        result = await sync_drive_folder()
        if result.get("error"):
            await send_telegram_message(chat_id, f"❌ {result['error']}")
        else:
            msg = f"✅ 동기화 완료\n📥 신규: {result['indexed_count']}\n⏭ 스킵: {result['skipped_count']}\n❌ 실패: {result['failed_count']}"
            if result['indexed']:
                msg += "\n\n*신규:*\n" + "\n".join(f"• {n}" for n in result['indexed'][:10])
            await send_telegram_message(chat_id, msg)
        return {"ok": True}
    
    if text.strip() == "/help":
        await send_telegram_message(chat_id,
            "*기능*\n"
            "☀️ `/morning` — 통합 아침 브리핑\n"
            "📡 `/dm` 또는 `/dm 48` — 텔레그램 정리\n"
            "🧠 `/remember <내용>` — 영구 기억\n"
            "🧠 `/memories` — 기억 목록\n"
            "🧠 `/forget <키워드>` — 기억 삭제\n"
            "📅 캘린더 조회/추가\n"
            "💬 슬랙 조회/발송\n"
            "📚 문서 업로드 + RAG\n\n"
            "*자동*\n매일 07:00 KST 통합 브리핑\n중요 정보 자동 기억\n\n"
            "*모델*\n/think /sonnet /quick\n\n"
            "*기타*\n/reset /sync /auth /help")
        return {"ok": True}
    
    if text.strip() == "/auth":
        await send_telegram_message(chat_id, f"https://{RAILWAY_URL}/auth/google")
        return {"ok": True}
    
    model_to_use = select_model(text)
    if "opus" in model_to_use:
        await send_telegram_message(chat_id, "🧠 깊게 생각 중...")
    elif "sonnet" in model_to_use:
        await send_telegram_message(chat_id, "💭 분석 중...")
    else:
        await send_telegram_message(chat_id, "⚡ 처리 중...")
    
    reply = await get_claude_response(text)
    await send_telegram_message(chat_id, reply)
    return {"ok": True}


@app.post("/slack/events")
async def slack_events(request: Request):
    data = await request.json()
    if data.get("type") == "url_verification":
        return {"challenge": data.get("challenge")}
    return {"ok": True}


# ==================== Google OAuth ====================

_oauth_flow_instance = None


def _create_google_flow():
    return Flow.from_client_config({
        "web": {
            "client_id": GOOGLE_CLIENT_ID,
            "client_secret": GOOGLE_CLIENT_SECRET,
            "auth_uri": "https://accounts.google.com/o/oauth2/auth",
            "token_uri": "https://oauth2.googleapis.com/token",
            "redirect_uris": [REDIRECT_URI]
        }
    }, scopes=GOOGLE_SCOPES)


@app.get("/auth/google")
def auth_google():
    global _oauth_flow_instance
    _oauth_flow_instance = _create_google_flow()
    _oauth_flow_instance.redirect_uri = REDIRECT_URI
    auth_url, _ = _oauth_flow_instance.authorization_url(
        access_type='offline', prompt='consent', include_granted_scopes='true'
    )
    return RedirectResponse(auth_url)


@app.get("/auth/google/callback")
def auth_google_callback(code: str):
    global _oauth_flow_instance
    if _oauth_flow_instance is None:
        return HTMLResponse("<h1>먼저 /auth/google 방문</h1>")
    _oauth_flow_instance.fetch_token(code=code)
    refresh_token = _oauth_flow_instance.credentials.refresh_token
    _oauth_flow_instance = None
    return HTMLResponse(f"""<html><body style="font-family:sans-serif;padding:40px"><h1>✅ 인증 성공</h1><p>Railway GOOGLE_REFRESH_TOKEN을 아래 값으로 업데이트:</p><pre style="background:#f0f0f0;padding:20px;border-radius:8px;word-break:break-all">{refresh_token}</pre></body></html>""")


@app.get("/")
def root():
    mem_count = memory_stats().get('count', 0)
    return {
        "status": "running",
        "user": USER_NAME,
        "company": COMPANY_NAME,
        "google": bool(GOOGLE_REFRESH_TOKEN),
        "slack": bool(SLACK_TOKEN),
        "pinecone": bool(pinecone_index),
        "voyage": bool(voyage_client),
        "drive_folder": bool(GOOGLE_DRIVE_KB_FOLDER_ID),
        "telethon": bool(TELETHON_SESSION),
        "scheduler_active": _scheduler is not None and _scheduler.running if _scheduler else False,
        "memory_count": mem_count,
        "dm_scanning": _dm_scan_running,
        "morning_briefing_running": _morning_briefing_running,
        "history": len(conversation_history)
    }
